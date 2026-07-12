import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from jinja2 import Environment, FileSystemLoader
import urllib.parse

# 설정값
PAGES_URL = "https://meesokim.github.io/my-youtube-gallary"
OUTPUT_DIR = "docs"

def create_pptx(video_data, output_path):
    """ 각 페이지당 하나의 핵심 인사이트만 명확히 담는 다중 슬라이드 PPTX 생성 """
    prs = Presentation()
    # 16:9 와이드스크린 설정
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    # 헬퍼 함수: 슬라이드 추가 및 배경색(네이비) 설정
    def add_custom_slide():
        slide = prs.slides.add_slide(blank_layout)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0x0f, 0x0f, 0x1a) # 어두운 네이비
        return slide

    # ----------------------------------------------------
    # SLIDE 1: Cover & Overview (표지 및 전체 요약)
    # ----------------------------------------------------
    slide1 = add_custom_slide()
    
    # 카테고리 라벨
    cat_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.333), Inches(0.5))
    p_cat = cat_box.text_frame.paragraphs[0]
    p_cat.text = f"💡 [ {video_data['category']} ]"
    p_cat.font.name = "Malgun Gothic"
    p_cat.font.size = Pt(22)
    p_cat.font.bold = True
    p_cat.font.color.rgb = RGBColor(0xff, 0x6b, 0x6b)
    
    # 대제목
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.333), Inches(2.0))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = video_data['title']
    p_title.font.name = "Malgun Gothic"
    p_title.font.size = Pt(46)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    
    # 구분선
    line_shape = slide1.shapes.add_shape(1, Inches(1.0), Inches(4.0), Inches(2.5), Inches(0.06))
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = RGBColor(0x66, 0x7e, 0xea)
    line_shape.line.fill.background()
    
    # 개요 (Overview)
    ov_box = slide1.shapes.add_textbox(Inches(1.0), Inches(4.4), Inches(11.333), Inches(2.0))
    tf_ov = ov_box.text_frame
    tf_ov.word_wrap = True
    p_ov = tf_ov.paragraphs[0]
    p_ov.text = video_data['overview']
    p_ov.font.name = "Malgun Gothic"
    p_ov.font.size = Pt(24)
    p_ov.font.color.rgb = RGBColor(0xa0, 0xa0, 0xb8)
    p_ov.line_spacing = 1.3

    # ----------------------------------------------------
    # SLIDE 2, 3, 4: Core Insights (각 장표당 하나의 명확한 인사이트)
    # ----------------------------------------------------
    for idx, point_str in enumerate(video_data['points'][:3]):
        slide = add_custom_slide()
        
        # 상단 네비게이터 (현재 위치 파악용)
        nav_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(11.333), Inches(0.4))
        p_nav = nav_box.text_frame.paragraphs[0]
        p_nav.text = f"CORE INSIGHT 0{idx + 1} / 03"
        p_nav.font.name = "Malgun Gothic"
        p_nav.font.size = Pt(12)
        p_nav.font.bold = True
        p_nav.font.color.rgb = RGBColor(0x66, 0x7e, 0xea)
        
        # '제목: 상세설명' 파싱
        insight_title = f"Insight {idx + 1}"
        insight_detail = point_str
        if ":" in point_str:
            parts = point_str.split(":", 1)
            insight_title = parts[0].strip()
            insight_detail = parts[1].strip()
            
        # 인사이트 제목 (큼직하고 굵게)
        ins_title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(1.2))
        tf_ins_title = ins_title_box.text_frame
        tf_ins_title.word_wrap = True
        p_ins_title = tf_ins_title.paragraphs[0]
        p_ins_title.text = insight_title
        p_ins_title.font.name = "Malgun Gothic"
        p_ins_title.font.size = Pt(30)
        p_ins_title.font.bold = True
        p_ins_title.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
        
        # 구분선
        line_shape = slide.shapes.add_shape(1, Inches(1.0), Inches(2.5), Inches(11.333), Inches(0.02))
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = RGBColor(0x2d, 0x2d, 0x44)
        line_shape.line.fill.background()
        
        # 인사이트 상세 내용 (1 Slide, 1 Message를 위한 큼직한 가독성 확보)
        ins_detail_box = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.333), Inches(3.5))
        tf_ins_detail = ins_detail_box.text_frame
        tf_ins_detail.word_wrap = True
        p_ins_detail = tf_ins_detail.paragraphs[0]
        p_ins_detail.text = insight_detail
        p_ins_detail.font.name = "Malgun Gothic"
        p_ins_detail.font.size = Pt(20)
        p_ins_detail.font.color.rgb = RGBColor(0xe8, 0xe8, 0xf0)
        p_ins_detail.line_spacing = 1.4

    # ----------------------------------------------------
    # SLIDE 5: Outro & Recommendation (마무리 및 시청자 추천사)
    # ----------------------------------------------------
    slide5 = add_custom_slide()
    
    # 추천사 카드 형태의 백그라운드 영역
    card_shape = slide5.shapes.add_shape(
        5, Inches(1.0), Inches(1.5), Inches(11.333), Inches(3.8) # 5 = rounded rectangle
    )
    card_shape.fill.solid()
    card_shape.fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
    card_shape.line.color.rgb = RGBColor(0x2d, 0x2d, 0x44)
    
    # 추천사 텍스트
    rec_box = slide5.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.333), Inches(3.2))
    tf_rec = rec_box.text_frame
    tf_rec.word_wrap = True
    
    p_rec_lbl = tf_rec.paragraphs[0]
    p_rec_lbl.text = "🎯 AI 추천사 및 시청 가치"
    p_rec_lbl.font.name = "Malgun Gothic"
    p_rec_lbl.font.size = Pt(14)
    p_rec_lbl.font.bold = True
    p_rec_lbl.font.color.rgb = RGBColor(0xff, 0x6b, 0x6b)
    p_rec_lbl.space_after = Pt(14)
    
    p_rec_val = tf_rec.add_paragraph()
    p_rec_val.text = video_data.get('recommendation', '영상을 통해 최신 통찰을 넓혀보세요.')
    p_rec_val.font.name = "Malgun Gothic"
    p_rec_val.font.size = Pt(18)
    p_rec_val.font.italic = True
    p_rec_val.font.color.rgb = RGBColor(0xe8, 0xe8, 0xf0)
    p_rec_val.line_spacing = 1.3
    
    # 하단 메타바 (키워드 및 출처 링크)
    meta_box = slide5.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(11.333), Inches(1.0))
    tf_meta = meta_box.text_frame
    tf_meta.word_wrap = True
    p_meta = tf_meta.paragraphs[0]
    p_meta.text = f"🏷️ Keywords: {video_data['keywords']}\n🔗 Source: https://youtu.be/{video_data['youtube_id']}"
    p_meta.font.name = "Malgun Gothic"
    p_meta.font.size = Pt(12)
    p_meta.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
    p_meta.line_spacing = 1.3
    
    prs.save(output_path)

def build_site():
    # 출력 디렉토리 생성
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    
    # 기존 빌드 결과물 정리 (기존에 있던 .html, .pptx 파일만 지워줌)
    for filename in os.listdir(OUTPUT_DIR):
        if filename.endswith(".html") or filename.endswith(".pptx"):
            file_path = os.path.join(OUTPUT_DIR, filename)
            try:
                os.remove(file_path)
                print(f"[-] 기존 파일 제거: {file_path}")
            except Exception as e:
                print(f"[!] 기존 파일 제거 중 오류 발생: {e}")
                
    # 데이터 로드
    with open('data/videos.json', 'r', encoding='utf-8') as f:
        videos = json.load(f)
        
    env = Environment(loader=FileSystemLoader('templates'))
    
    # URL 인코딩을 위한 헬퍼 함수 등록
    env.globals['encode_url'] = urllib.parse.quote
    env.globals['pages_url'] = PAGES_URL
    
    # 1. 상세 페이지 및 PPTX 생성
    preview_template = env.get_template('video.html')
    for video in videos:
        filename = video['filename']
        
        # PPTX 저장
        pptx_path = os.path.join(OUTPUT_DIR, f"{filename}.pptx")
        create_pptx(video, pptx_path)
        print(f"[+] PPTX 생성: {pptx_path}")
        
        # 상세 HTML 저장
        rendered_html = preview_template.render(video=video)
        html_path = os.path.join(OUTPUT_DIR, f"{filename}.html")
        with open(html_path, "w", encoding="utf-8") as f:
            f.write(rendered_html)
        print(f"[+] HTML 생성: {html_path}")
            
    # 2. 메인 index.html 생성
    index_template = env.get_template('index.html')
    rendered_index = index_template.render(videos=videos)
    index_path = os.path.join(OUTPUT_DIR, "index.html")
    with open(index_path, "w", encoding="utf-8") as f:
        f.write(rendered_index)
    print(f"[+] Index 생성: {index_path}")
    print(f"[*] 총 {len(videos)}개 영상 처리 완료!")

if __name__ == "__main__":
    build_site()
